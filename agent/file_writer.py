"""
Rich file writer for coding-agent v2.0.2

Supports writing AI-generated content to:
  .docx  — Word document   (python-docx)
  .pptx  — PowerPoint      (python-pptx)
  .xlsx  — Excel workbook  (openpyxl)
  .pdf   — PDF document    (fpdf2)
  .md    — Markdown        (plain text)
  .txt   — Plain text      (plain text)
  .csv   — CSV table       (csv module)

Flow:
  1. User gives a natural-language instruction
     e.g. "Create a project summary for WeatherAgent with architecture overview"
  2. AI (Ollama) generates a structured JSON document plan
  3. This module converts the JSON plan → actual binary/text file

JSON Schema the AI must produce:
─────────────────────────────────
For .docx / .pdf / .md / .txt:
{
  "title": "string",
  "subtitle": "string (optional)",
  "sections": [
    {
      "heading": "string",
      "level": 1,                        // 1=H1, 2=H2, 3=H3
      "paragraphs": ["string"],          // optional
      "bullets": ["string"],             // optional
      "table": {                         // optional
        "headers": ["col1", "col2"],
        "rows": [["val1", "val2"]]
      },
      "code": {                          // optional
        "language": "python",
        "content": "string"
      }
    }
  ]
}

For .pptx:
{
  "title": "string",
  "subtitle": "string (optional)",
  "slides": [
    {
      "title": "string",
      "layout": "title | bullets | two_col | table | blank",
      "content": "string (optional)",
      "bullets": ["string"],
      "table": {
        "headers": ["col1"],
        "rows": [["val1"]]
      },
      "notes": "string (optional speaker notes)"
    }
  ]
}

For .xlsx:
{
  "title": "string",
  "sheets": [
    {
      "name": "string",
      "headers": ["col1", "col2"],
      "rows": [["val1", "val2"]],
      "summary": "string (optional — written above the table)"
    }
  ]
}
"""

from __future__ import annotations

import csv
import io
import json
import re
import textwrap
from pathlib import Path
from typing import Any


# ═══════════════════════════════════════════════════════════════════════════════
# AI Prompt — generates the JSON doc plan
# ═══════════════════════════════════════════════════════════════════════════════

_DOC_SYSTEM = textwrap.dedent("""\
    You are a technical writer and document architect.
    Given a user instruction you produce a structured JSON document plan.
    Rules:
    - Return ONLY a valid JSON object. No markdown fences. No explanation outside JSON.
    - Follow the schema exactly as described for the target file type.
    - Be thorough: include real content, not placeholders.
    - Tables must have at least 2 rows of actual data.
""")


def _doc_prompt(instruction: str, fmt: str, context: str = "") -> str:
    schema_hint = {
        "docx": 'Schema: {"title":…,"sections":[{"heading":…,"level":1,"paragraphs":[…],"bullets":[…],"table":{"headers":[…],"rows":[[…]]}}]}',
        "pdf":  'Schema: {"title":…,"sections":[{"heading":…,"level":1,"paragraphs":[…],"bullets":[…],"table":{"headers":[…],"rows":[[…]]}}]}',
        "md":   'Schema: {"title":…,"sections":[{"heading":…,"level":1,"paragraphs":[…],"bullets":[…],"table":{"headers":[…],"rows":[[…]]}}]}',
        "txt":  'Schema: {"title":…,"sections":[{"heading":…,"level":1,"paragraphs":[…],"bullets":[…]}]}',
        "pptx": 'Schema: {"title":…,"subtitle":…,"slides":[{"title":…,"layout":"bullets","bullets":[…],"notes":…}]}',
        "xlsx": 'Schema: {"title":…,"sheets":[{"name":…,"headers":[…],"rows":[[…]],"summary":…}]}',
        "csv":  'Schema: {"title":…,"sheets":[{"name":"Sheet1","headers":[…],"rows":[[…]]}]}',
    }.get(fmt, "")

    ctx_block = f"\n\nProject context:\n{context[:3000]}" if context else ""

    return (
        f"Create a {fmt.upper()} document for the following instruction.\n\n"
        f"Instruction: {instruction}\n\n"
        f"{schema_hint}{ctx_block}\n\n"
        "Return ONLY the JSON object."
    )


def _extract_json(text: str) -> str:
    if not text or not isinstance(text, str):
        return "{}"
    text = re.sub(r"```(?:json)?\s*", "", text)
    text = re.sub(r"```", "", text)
    start = text.find("{")
    if start == -1:
        return text.strip()
    depth = 0
    for i, ch in enumerate(text[start:], start):
        if ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                return text[start: i + 1]
    return text[start:].strip()


def _fallback_pptx_plan(instruction: str, context: str) -> dict:
    """
    Build a minimal valid pptx plan from the context text when the LLM
    fails to produce valid JSON. Splits the context into sections by
    heading lines (## / ###) and maps each to a slide.
    """
    lines = (context or instruction).splitlines()
    title = "Presentation"
    subtitle = ""
    slides: list[dict] = []
    current_title = ""
    current_bullets: list[str] = []

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            title = line.lstrip("# ").strip()
        elif line.startswith("## ") or line.startswith("### "):
            if current_title:
                slides.append({"title": current_title, "layout": "bullets",
                                "bullets": current_bullets or [""], "notes": ""})
            current_title = line.lstrip("# ").strip()
            current_bullets = []
        elif line.startswith("- ") or line.startswith("* "):
            current_bullets.append(line.lstrip("-* ").strip())
        elif current_title and not line.startswith("#"):
            current_bullets.append(line)

    if current_title:
        slides.append({"title": current_title, "layout": "bullets",
                        "bullets": current_bullets or [""], "notes": ""})

    if not slides:
        # Absolute fallback — one slide per paragraph
        for i, chunk in enumerate(lines[:10]):
            if chunk:
                slides.append({"title": f"Slide {i+1}", "layout": "bullets",
                                "bullets": [chunk], "notes": ""})

    return {"title": title, "subtitle": subtitle, "slides": slides[:15]}


def plan_document(instruction: str, fmt: str, engine: Any, context: str = "") -> dict:
    """
    Ask Ollama to generate a document plan as JSON.
    Retries once with a simpler prompt, then falls back to a
    structure derived directly from the source context.
    """
    for attempt in range(2):
        prompt = _doc_prompt(instruction, fmt, context) if attempt == 0 else (
            f"Return ONLY a JSON object for a {fmt.upper()} with this schema:\n"
            + {
                "pptx": '{"title":"…","subtitle":"…","slides":[{"title":"…","layout":"bullets","bullets":["…"],"notes":""}]}',
                "docx": '{"title":"…","sections":[{"heading":"…","level":1,"paragraphs":["…"],"bullets":["…"]}]}',
                "xlsx": '{"title":"…","sheets":[{"name":"Sheet1","headers":["…"],"rows":[["…"]],"summary":"…"}]}',
                "pdf":  '{"title":"…","sections":[{"heading":"…","level":1,"paragraphs":["…"],"bullets":["…"]}]}',
            }.get(fmt, '{"title":"…","sections":[]}')
            + f"\n\nContent to use:\n{(context or instruction)[:2000]}"
            + "\n\nReturn ONLY valid JSON. No explanation."
        )
        raw = engine.complete(prompt, system=_DOC_SYSTEM, temperature=0.1)
        json_str = _extract_json(raw)
        try:
            plan = json.loads(json_str)
            if plan:
                return plan
        except (json.JSONDecodeError, ValueError):
            pass

    # Both LLM attempts failed — build plan directly from context
    if fmt == "pptx":
        return _fallback_pptx_plan(instruction, context)
    # Generic fallback for other formats
    return {
        "title": instruction[:60],
        "sections": [{"heading": "Content", "level": 1,
                      "paragraphs": [(context or instruction)[:500]], "bullets": []}],
    }


# ═══════════════════════════════════════════════════════════════════════════════
# Writers
# ═══════════════════════════════════════════════════════════════════════════════

def write_docx(plan: dict, output: Path) -> None:
    """Write a Word .docx file from a document plan."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    def shade_hdr(cell, hex_color="4472C4"):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

    doc = Document()

    # Title
    title_para = doc.add_heading(plan.get("title", "Document"), 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if plan.get("subtitle"):
        sub = doc.add_paragraph(plan["subtitle"])
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.color.rgb = RGBColor(100, 100, 100)
        sub.runs[0].font.size = Pt(13)

    doc.add_paragraph()

    for section in plan.get("sections", []):
        level = min(max(section.get("level", 1), 1), 3)
        heading = section.get("heading", "")
        if heading:
            doc.add_heading(heading, level)

        for para in section.get("paragraphs", []):
            if para.strip():
                doc.add_paragraph(para)

        for bullet in section.get("bullets", []):
            if bullet.strip():
                doc.add_paragraph(bullet, style="List Bullet")

        code_block = section.get("code")
        if code_block:
            p = doc.add_paragraph()
            run = p.add_run(code_block.get("content", ""))
            run.font.name = "Courier New"
            run.font.size = Pt(9)

        tbl_data = section.get("table")
        if tbl_data:
            headers = tbl_data.get("headers", [])
            rows    = tbl_data.get("rows", [])
            if headers:
                tbl = doc.add_table(rows=1, cols=len(headers))
                tbl.style = "Table Grid"
                hdr_cells = tbl.rows[0].cells
                for i, h in enumerate(headers):
                    hdr_cells[i].text = str(h)
                    shade_hdr(hdr_cells[i])
                    hdr_cells[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    hdr_cells[i].paragraphs[0].runs[0].bold = True
                for row_data in rows:
                    row_cells = tbl.add_row().cells
                    for i, val in enumerate(row_data):
                        if i < len(row_cells):
                            row_cells[i].text = str(val)
                doc.add_paragraph()

    output.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output))


def write_pptx(plan: dict, output: Path) -> None:
    """Write a PowerPoint .pptx from a document plan."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    LAYOUTS = {name: i for i, name in enumerate([
        "Title Slide", "Title and Content", "Section Header",
        "Two Content", "Comparison", "Title Only",
        "Blank", "Content with Caption", "Picture with Caption",
        "Title and Vertical Text", "Vertical Title and Text",
    ])}

    def get_layout(name: str):
        idx = {"title": 0, "bullets": 1, "section": 2,
               "two_col": 3, "table": 5, "blank": 6}.get(name, 1)
        return prs.slide_layouts[min(idx, len(prs.slide_layouts) - 1)]

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = plan.get("title", "Presentation")
    if plan.get("subtitle") and title_slide.placeholders[1]:
        title_slide.placeholders[1].text = plan.get("subtitle", "")

    for slide_data in plan.get("slides", []):
        layout_name = slide_data.get("layout", "bullets")
        slide = prs.slides.add_slide(get_layout(layout_name))

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")

        # Content / bullets
        bullets = slide_data.get("bullets", [])
        content = slide_data.get("content", "")

        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break

        if body_ph:
            tf = body_ph.text_frame
            tf.clear()
            if content:
                tf.text = content
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        # Table
        tbl_data = slide_data.get("table")
        if tbl_data:
            headers = tbl_data.get("headers", [])
            rows    = tbl_data.get("rows", [])
            if headers and rows:
                n_cols = len(headers)
                n_rows = len(rows) + 1
                left   = Inches(1.0)
                top    = Inches(2.5)
                width  = Inches(11.0)
                height = Inches(0.4 * n_rows)
                tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
                for c, h in enumerate(headers):
                    cell = tbl.cell(0, c)
                    cell.text = str(h)
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                for r, row_data in enumerate(rows):
                    for c, val in enumerate(row_data):
                        if c < n_cols:
                            tbl.cell(r + 1, c).text = str(val)

        # Speaker notes
        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def write_xlsx(plan: dict, output: Path) -> None:
    """Write an Excel .xlsx workbook from a document plan."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    wb.remove(wb.active)   # remove default empty sheet

    hdr_fill   = PatternFill("solid", fgColor="4472C4")
    hdr_font   = Font(bold=True, color="FFFFFF", size=11)
    alt_fill   = PatternFill("solid", fgColor="D9E1F2")
    thin_side  = Side(style="thin", color="AAAAAA")
    thin_border = Border(left=thin_side, right=thin_side,
                         top=thin_side, bottom=thin_side)

    for sheet_data in plan.get("sheets", []):
        ws = wb.create_sheet(title=sheet_data.get("name", "Sheet")[:31])
        headers = sheet_data.get("headers", [])
        rows    = sheet_data.get("rows", [])
        summary = sheet_data.get("summary", "")

        start_row = 1
        if summary:
            ws.cell(row=1, column=1, value=summary).font = Font(italic=True, color="555555")
            start_row = 3

        # Headers
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=start_row, column=c, value=str(h))
            cell.fill   = hdr_fill
            cell.font   = hdr_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

        # Data rows
        for r, row_data in enumerate(rows, start_row + 1):
            fill = alt_fill if r % 2 == 0 else PatternFill()
            for c, val in enumerate(row_data, 1):
                if c <= len(headers):
                    cell = ws.cell(row=r, column=c, value=val)
                    cell.fill   = fill
                    cell.border = thin_border
                    cell.alignment = Alignment(vertical="center")

        # Auto-size columns
        for c in range(1, len(headers) + 1):
            col_letter = get_column_letter(c)
            max_len = max(
                (len(str(ws.cell(row=r, column=c).value or ""))
                 for r in range(start_row, start_row + len(rows) + 1)),
                default=10,
            )
            ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

    output.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output))


def write_pdf(plan: dict, output: Path) -> None:
    """Write a PDF from a document plan using fpdf2."""
    from fpdf import FPDF

    class PDF(FPDF):
        def header(self):
            pass
        def footer(self):
            self.set_y(-12)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(150, 150, 150)
            self.cell(0, 8, f"Page {self.page_no()}", align="C")

    pdf = PDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()
    pdf.set_margins(18, 18, 18)

    # Title
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(31, 56, 100)
    pdf.multi_cell(0, 12, plan.get("title", "Document"), align="C")
    pdf.ln(3)

    if plan.get("subtitle"):
        pdf.set_font("Helvetica", "I", 13)
        pdf.set_text_color(100, 100, 100)
        pdf.multi_cell(0, 8, plan["subtitle"], align="C")
        pdf.ln(6)

    pdf.set_draw_color(68, 114, 196)
    pdf.set_line_width(0.5)
    pdf.line(18, pdf.get_y(), 192, pdf.get_y())
    pdf.ln(8)

    LEVEL_SIZES  = {1: 16, 2: 13, 3: 11}
    LEVEL_COLORS = {1: (31, 56, 100), 2: (68, 114, 196), 3: (80, 80, 80)}

    for section in plan.get("sections", []):
        level   = min(max(section.get("level", 1), 1), 3)
        heading = section.get("heading", "")

        if heading:
            pdf.set_font("Helvetica", "B", LEVEL_SIZES[level])
            r, g, b = LEVEL_COLORS[level]
            pdf.set_text_color(r, g, b)
            pdf.multi_cell(0, 9, heading)
            pdf.ln(2)

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(40, 40, 40)

        for para in section.get("paragraphs", []):
            if para.strip():
                pdf.multi_cell(0, 6, para)
                pdf.ln(3)

        for bullet in section.get("bullets", []):
            if bullet.strip():
                pdf.set_x(24)
                pdf.cell(5, 6, chr(149))   # bullet char
                pdf.multi_cell(0, 6, bullet)
        if section.get("bullets"):
            pdf.ln(3)

        tbl = section.get("table")
        if tbl:
            headers = tbl.get("headers", [])
            rows    = tbl.get("rows", [])
            if headers:
                n_cols  = len(headers)
                col_w   = (pdf.w - 36) / n_cols

                # Header row
                pdf.set_font("Helvetica", "B", 9)
                pdf.set_fill_color(68, 114, 196)
                pdf.set_text_color(255, 255, 255)
                for h in headers:
                    pdf.cell(col_w, 8, str(h)[:30], border=1, fill=True, align="C")
                pdf.ln()

                # Data rows
                for r_idx, row_data in enumerate(rows):
                    pdf.set_font("Helvetica", "", 9)
                    fill = r_idx % 2 == 0
                    pdf.set_fill_color(217, 225, 242)
                    pdf.set_text_color(40, 40, 40)
                    for c, val in enumerate(row_data):
                        if c < n_cols:
                            pdf.cell(col_w, 7, str(val)[:35], border=1, fill=fill)
                    pdf.ln()
                pdf.ln(4)

        pdf.ln(4)

    output.parent.mkdir(parents=True, exist_ok=True)
    pdf.output(str(output))


def write_markdown(plan: dict, output: Path) -> None:
    """Write a Markdown .md file from a document plan."""
    lines: list[str] = []
    lines.append(f"# {plan.get('title', 'Document')}\n")

    if plan.get("subtitle"):
        lines.append(f"_{plan['subtitle']}_\n")

    lines.append("")

    for section in plan.get("sections", []):
        level   = min(max(section.get("level", 1), 1), 6)
        heading = section.get("heading", "")
        if heading:
            lines.append(f"{'#' * (level + 1)} {heading}\n")

        for para in section.get("paragraphs", []):
            if para.strip():
                lines.append(f"{para}\n")

        for bullet in section.get("bullets", []):
            if bullet.strip():
                lines.append(f"- {bullet}")
        if section.get("bullets"):
            lines.append("")

        code = section.get("code")
        if code:
            lang    = code.get("language", "")
            content = code.get("content", "")
            lines.append(f"```{lang}")
            lines.append(content)
            lines.append("```\n")

        tbl = section.get("table")
        if tbl:
            headers = tbl.get("headers", [])
            rows    = tbl.get("rows", [])
            if headers:
                lines.append("| " + " | ".join(str(h) for h in headers) + " |")
                lines.append("| " + " | ".join("---" for _ in headers) + " |")
                for row_data in rows:
                    lines.append("| " + " | ".join(str(v) for v in row_data) + " |")
                lines.append("")

        lines.append("")

    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")


def write_txt(plan: dict, output: Path) -> None:
    """Write a plain text .txt file from a document plan."""
    lines: list[str] = []
    title = plan.get("title", "Document")
    lines.append(title)
    lines.append("=" * len(title))
    lines.append("")

    if plan.get("subtitle"):
        lines.append(plan["subtitle"])
        lines.append("")

    for section in plan.get("sections", []):
        heading = section.get("heading", "")
        if heading:
            lines.append(heading)
            lines.append("-" * len(heading))

        for para in section.get("paragraphs", []):
            if para.strip():
                lines.append(para)
                lines.append("")

        for bullet in section.get("bullets", []):
            if bullet.strip():
                lines.append(f"  • {bullet}")
        if section.get("bullets"):
            lines.append("")

        tbl = section.get("table")
        if tbl:
            headers = tbl.get("headers", [])
            rows    = tbl.get("rows", [])
            if headers:
                widths = [max(len(str(h)), max((len(str(r[i])) for r in rows if i < len(r)), default=0))
                          for i, h in enumerate(headers)]
                sep  = "+-" + "-+-".join("-" * w for w in widths) + "-+"
                hrow = "| " + " | ".join(str(h).ljust(widths[i]) for i, h in enumerate(headers)) + " |"
                lines.extend([sep, hrow, sep])
                for row_data in rows:
                    drow = "| " + " | ".join(str(v).ljust(widths[i]) if i < len(widths) else ""
                                             for i, v in enumerate(row_data)) + " |"
                    lines.append(drow)
                lines.append(sep)
                lines.append("")

        lines.append("")

    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")


def write_csv(plan: dict, output: Path) -> None:
    """Write a CSV file — uses first sheet from the plan."""
    sheets = plan.get("sheets", [])
    if not sheets:
        output.write_text("", encoding="utf-8")
        return
    sheet   = sheets[0]
    headers = sheet.get("headers", [])
    rows    = sheet.get("rows", [])
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        if headers:
            writer.writerow(headers)
        writer.writerows(rows)


# ═══════════════════════════════════════════════════════════════════════════════
# Dispatcher
# ═══════════════════════════════════════════════════════════════════════════════

SUPPORTED_FORMATS = {".docx", ".pptx", ".xlsx", ".pdf", ".md", ".txt", ".csv"}

_WRITERS = {
    ".docx": write_docx,
    ".pptx": write_pptx,
    ".xlsx": write_xlsx,
    ".pdf":  write_pdf,
    ".md":   write_markdown,
    ".txt":  write_txt,
    ".csv":  write_csv,
}

# Formats that need the "sheets" schema instead of "sections"
_SHEET_FORMATS = {".xlsx", ".csv"}
# Formats that need the "slides" schema
_SLIDE_FORMATS = {".pptx"}


def write_document(
    instruction: str,
    output: Path,
    engine: Any,
    context: str = "",
) -> Path:
    """
    High-level entry point.
    Ask Ollama to plan the document, then write it to disk.

    Args:
        instruction: Natural language description of the document to create.
        output:      Destination path (extension determines format).
        engine:      CodingEngine instance (Ollama wrapper).
        context:     Optional extra context (e.g. project file contents).

    Returns:
        The resolved output path.
    """
    fmt = output.suffix.lower()
    if fmt not in SUPPORTED_FORMATS:
        raise ValueError(
            f"Unsupported format: {fmt}\n"
            f"Supported: {', '.join(sorted(SUPPORTED_FORMATS))}"
        )

    plan = plan_document(instruction, fmt.lstrip("."), engine, context)
    writer = _WRITERS[fmt]
    writer(plan, output)
    return output.resolve()
